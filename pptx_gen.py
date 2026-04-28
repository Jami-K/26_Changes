import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY   = RGBColor(0x1A, 0x35, 0x6E)
LIGHT  = RGBColor(0xE8, 0xEC, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x22, 0x22, 0x22)
BORDER = RGBColor(0xC0, 0xC0, 0xC0)

FIELDS = [
    ('notice_date',      '고지일'),
    ('product_code',     '제품코드'),
    ('product_name',     '제품명'),
    ('changed_material', '변경자재'),
    ('mixable',          '혼용가능여부'),
    ('change_content',   '변경내용'),
    ('apply_date',       '적용예정일'),
    ('notes',            '특이사항'),
]


def _rect(slide, left, top, w, h, fill, border=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    return shp


def _text(shp, content, size, bold, color, align=PP_ALIGN.LEFT):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left   = Pt(8)
    tf.margin_right  = Pt(8)
    tf.margin_top    = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = content
    run.font.name  = '맑은 고딕'
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color


def generate_pptx(change: dict) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # ── 헤더 배경 ──────────────────────────────────────────────
    hdr = _rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.15), NAVY)
    _text(hdr, '자재 변경점', Pt(30), True, WHITE, PP_ALIGN.CENTER)

    # ── 필드 테이블 ────────────────────────────────────────────
    x_start   = Inches(1.3)
    y_start   = Inches(1.35)
    lbl_width = Inches(1.85)
    val_width = Inches(9.68)   # 1.3 + 1.85 + 9.68 + 0.5 = 13.33
    row_h     = Inches(0.67)   # 8 rows × 0.67 = 5.36 → bottom at 6.71"

    for i, (key, label) in enumerate(FIELDS):
        top = y_start + i * row_h

        lbl_shp = _rect(slide, x_start, top, lbl_width, row_h, LIGHT, BORDER)
        _text(lbl_shp, label, Pt(12), True, NAVY, PP_ALIGN.CENTER)

        val_shp = _rect(slide, x_start + lbl_width, top, val_width, row_h, WHITE, BORDER)
        _text(val_shp, change.get(key) or '', Pt(12), False, DARK, PP_ALIGN.LEFT)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
